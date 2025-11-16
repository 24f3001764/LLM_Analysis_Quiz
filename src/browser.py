from playwright.async_api import async_playwright, Page, ElementHandle
from typing import Union, Optional, Dict, Any, List, Tuple, BinaryIO
import asyncio
import logging
from pathlib import Path
import json
import base64
from urllib.parse import urljoin, urlparse
import mimetypes
import magic
import re
import io
import tempfile
from datetime import datetime

# File processing libraries
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
    
try:
    import docx
except ImportError:
    docx = None
    
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    
try:
    import openpyxl
except ImportError:
    openpyxl = None
    
try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None

from .config import settings

logger = logging.getLogger(__name__)

# Supported file types and their MIME types
SUPPORTED_FILE_TYPES = {
    'application/pdf': 'pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
    'application/msword': 'doc',
    'application/vnd.ms-powerpoint': 'ppt',
    'application/vnd.ms-excel': 'xls',
    'text/plain': 'txt',
    'text/csv': 'csv',
    'image/png': 'png',
    'image/jpeg': 'jpg',
    'image/gif': 'gif',
    'image/tiff': 'tiff',
    'image/bmp': 'bmp'
}

# OCR supported image types
OCR_SUPPORTED = {'png', 'jpg', 'jpeg', 'tiff', 'bmp'}

class BrowserManager:
    """
    Manages browser instances and provides methods for web scraping.
    Uses Playwright for headless browser automation.
    """
    
    def __init__(self, headless: bool = True):
        self.headless = headless
        self.browser = None
        self.context = None
        self.page = None
        self._magic = magic.Magic(mime=True)
        self.playwright = None
    
    async def __aenter__(self):
        """Initialize the browser and context when entering the context manager."""
        self.playwright = await async_playwright().start()
        self.browser = await self.playwright.chromium.launch(
            headless=self.headless,
            args=[
                '--disable-blink-features=AutomationControlled',
                '--no-sandbox',
                '--disable-setuid-sandbox',
                '--disable-dev-shm-usage',
                '--disable-accelerated-2d-canvas',
                '--no-first-run',
                '--no-zygote',
                '--single-process',
                '--disable-gpu',
                '--disable-software-rasterizer',
                '--disable-features=site-per-process',
                '--disable-features=IsolateOrigins,site-per-process'
            ]
        )
        
        # Create a new browser context
        self.context = await self.browser.new_context(
            viewport={'width': 1920, 'height': 1080},
            user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
            java_script_enabled=True,
            ignore_https_errors=True,
            bypass_csp=True
        )
        
        # Add stealth mode to avoid detection
        await self.context.add_init_script("""
        // Overwrite the `languages` property to use a custom getter.
        Object.defineProperty(navigator, 'languages', {
            get: function() {
                return ['en-US', 'en'];
            },
        });
        
        // Overwrite the `plugins` property to use a custom getter.
        Object.defineProperty(navigator, 'plugins', {
            get: function() {
                // This just needs to have `length > 0`, but we could mock the plugins too
                return [1, 2, 3, 4, 5];
            },
        });
        
        // Overwrite the `plugins` property to use a custom getter.
        Object.defineProperty(navigator, 'webdriver', {
            get: function () {
                return false;
            },
        });
        """)
        
        # Create a new page
        self.page = await self.context.new_page()
        return self
    
    async def __aexit__(self, exc_type, exc_val, exc_tb):
        """Clean up resources when exiting the context manager."""
        if self.browser:
            await self.browser.close()
        if self.playwright:
            await self.playwright.stop()
    
    async def close(self):
        """Close the browser and all its contexts."""
        if self.browser:
            await self.browser.close()
            self.browser = None
            self.context = None
            self.page = None
            
    def _get_file_extension(self, file_path: Union[str, Path], content_type: str = None) -> str:
        """Get the file extension based on content type or file path."""
        if content_type:
            return SUPPORTED_FILE_TYPES.get(content_type.lower(), '')
        return Path(file_path).suffix.lstrip('.').lower()
    
    def _is_ocr_supported(self, file_extension: str) -> bool:
        """Check if the file type supports OCR."""
        return file_extension.lower() in OCR_SUPPORTED
    
    def _extract_text_from_pdf(self, file_path: Union[str, Path]) -> str:
        """Extract text from PDF file."""
        if not PyPDF2:
            raise ImportError("PyPDF2 is required for PDF processing. Install with: pip install PyPDF2")
            
        text = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text() or '')
        return '\n'.join(text)
    
    def _extract_text_from_docx(self, file_path: Union[str, Path]) -> str:
        """Extract text from DOCX file."""
        if not docx:
            raise ImportError("python-docx is required for DOCX processing. Install with: pip install python-docx")
            
        doc = docx.Document(file_path)
        return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
    
    def _extract_text_from_pptx(self, file_path: Union[str, Path]) -> str:
        """Extract text from PPTX file."""
        if not Presentation:
            raise ImportError("python-pptx is required for PPTX processing. Install with: pip install python-pptx")
            
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    
    def _extract_text_from_xlsx(self, file_path: Union[str, Path]) -> str:
        """Extract text from XLSX file."""
        if not openpyxl:
            raise ImportError("openpyxl is required for XLSX processing. Install with: pip install openpyxl")
            
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text.append(f"--- Sheet: {sheet_name} ---")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        if isinstance(cell, str):
                            row_text.append(cell.strip())
                        else:
                            row_text.append(str(cell))
                if any(row_text):  # Only add non-empty rows
                    text.append('\t'.join(row_text))
        
        return '\n'.join(text)
        
    def _extract_text_with_ocr(self, file_path: Union[str, Path]) -> str:
        """Extract text from image using OCR."""
        if not Image or not pytesseract:
            raise ImportError("Pillow and pytesseract are required for OCR. Install with: pip install Pillow pytesseract")
            
        try:
            # Open the image file
            with Image.open(file_path) as img:
                # Convert to grayscale for better OCR
                if img.mode != 'L':
                    img = img.convert('L')
                
                # Extract text using pytesseract
                text = pytesseract.image_to_string(
                    img,
                    config='--psm 6'  # Assume a single uniform block of text
                )
                
                return text.strip()
                
        except Exception as e:
            logger.error(f"OCR processing failed: {str(e)}")
            return ""
            
    def extract_text_from_file(self, file_path: Union[str, Path], content_type: str = None) -> Dict[str, Any]:
        """
        Extract text from various file types.
        
        Args:
            file_path: Path to the file
            content_type: Optional MIME type of the file
            
        Returns:
            Dict containing:
            - text: Extracted text
            - metadata: File metadata
            - content_type: Detected content type
            - file_size: File size in bytes
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        # Get file info
        file_size = file_path.stat().st_size
        file_name = file_path.name
        
        try:
            # Detect content type if not provided
            if not content_type:
                content_type = self._magic.from_file(str(file_path), mime=True)
            
            # Get file extension
            file_extension = self._get_file_extension(file_path, content_type)
            
            # Initialize result with metadata
            result = {
                'text': '',
                'metadata': {
                    'file_name': file_name,
                    'file_size': file_size,
                    'content_type': content_type,
                    'file_extension': file_extension,
                    'processing_time': datetime.utcnow().isoformat()
                },
                'content_type': content_type,
                'file_size': file_size
            }
            
            # Handle different file types
            try:
                if 'pdf' in content_type.lower():
                    result['text'] = self._extract_text_from_pdf(file_path)
                elif 'wordprocessingml' in content_type.lower() or 'msword' in content_type.lower():
                    result['text'] = self._extract_text_from_docx(file_path)
                elif 'presentationml' in content_type.lower() or 'powerpoint' in content_type.lower():
                    result['text'] = self._extract_text_from_pptx(file_path)
                elif 'spreadsheetml' in content_type.lower() or 'excel' in content_type.lower():
                    result['text'] = self._extract_text_from_xlsx(file_path)
                elif content_type.startswith('text/'):
                    result['text'] = file_path.read_text(encoding='utf-8', errors='replace')
                elif any(img_type in content_type.lower() for img_type in ['image/png', 'image/jpeg', 'image/tiff', 'image/bmp']):
                    result['text'] = self._extract_text_with_ocr(file_path)
                    result['metadata']['ocr_used'] = True
                else:
                    logger.warning(f"Unsupported content type for text extraction: {content_type}")
                    result['error'] = f"Unsupported content type: {content_type}"
                    
            except Exception as e:
                error_msg = f"Error extracting text from {file_name}: {str(e)}"
                logger.error(error_msg)
                result['error'] = error_msg
                
                # Try OCR as fallback for images
                if self._is_ocr_supported(file_extension) and 'ocr_used' not in result.get('metadata', {}):
                    try:
                        result['text'] = self._extract_text_with_ocr(file_path)
                        result['metadata']['ocr_used'] = True
                        result['metadata']['ocr_fallback'] = True
                    except Exception as ocr_error:
                        logger.error(f"OCR fallback failed: {str(ocr_error)}")
                        result['ocr_error'] = str(ocr_error)
            
            return result
            
        except Exception as e:
            error_msg = f"Failed to process file {file_name}: {str(e)}"
            logger.error(error_msg)
            return {
                'text': '',
                'error': error_msg,
                'metadata': {
                    'file_name': file_name,
                    'file_size': file_size,
                    'content_type': content_type or 'unknown',
                    'error': str(e)
                }
            }
    
    async def upload_file(self, file_input_selector: str, file_path: Union[str, Path]) -> None:
        """
        Upload a file to a file input element.
        
        Args:
            file_input_selector: CSS selector for the file input element
            file_path: Path to the file to upload
            
        Raises:
            FileNotFoundError: If the specified file does not exist
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        # Convert to absolute path
        file_path = file_path.absolute()
        
        # Set the file input value
        await self.page.set_input_files(file_input_selector, str(file_path))
        logger.info(f"Uploaded file: {file_path}")
    
    async def download_file(self, url: str, save_path: Optional[Path] = None, extract_text: bool = True) -> Dict[str, Any]:
        """
        Download a file from a given URL.
        
        Args:
            url: URL of the file to download
            save_path: Optional path to save the file. If None, uses settings.DOWNLOADS_DIR
            extract_text: Whether to extract text from the downloaded file
            
        Returns:
            Dict containing file metadata and content
            
        Example:
            ```python
            async with BrowserManager() as browser:
                result = await browser.download_file(
                    "https://example.com/document.pdf",
                    "downloads/document.pdf"
                )
                print(f"Downloaded {result['file_name']} ({result['file_size']} bytes)")
                if 'text' in result:
                    print(f"Extracted text: {result['text'][:200]}...")
            ```
        """
        try:
            if save_path is None:
                save_path = Path(settings.DOWNLOADS_DIR) / Path(url).name
            else:
                save_path = Path(save_path)
                
            save_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Start waiting for the download
            async with self.page.expect_download() as download_info:
                await self.page.goto(url)
            
            download = await download_info.value
            
            # Save the file
            await download.save_as(str(save_path))
            
            # Get file metadata
            content_type = self._magic.from_file(str(save_path), mime=True)
            file_size = save_path.stat().st_size
            
            result = {
                'url': url,
                'path': str(save_path),
                'file_name': save_path.name,
                'content_type': content_type,
                'file_size': file_size,
                'downloaded_at': datetime.utcnow().isoformat()
            }
            
            # Extract text if requested
            if extract_text:
                try:
                    text_result = self.extract_text_from_file(save_path, content_type)
                    result.update({
                        'text': text_result.get('text', ''),
                        'metadata': text_result.get('metadata', {})
                    })
                    if 'error' in text_result:
                        result['error'] = text_result['error']
                except Exception as e:
                    logger.error(f"Error extracting text from downloaded file: {str(e)}")
                    result['error'] = str(e)
            
            return result
            
        except Exception as e:
            error_msg = f"Error downloading file from {url}: {str(e)}"
            logger.error(error_msg)
            return {
                'url': url,
                'error': error_msg,
                'metadata': {
                    'error': str(e),
                    'timestamp': datetime.utcnow().isoformat()
                }
            }
    
    def extract_text_from_file(self, file_path: Union[str, Path], content_type: str = None) -> Dict[str, Any]:
        """
        Extract text from various file types.
        
        Args:
            file_path: Path to the file
            content_type: Optional MIME type of the file
            
        Returns:
            Dict containing:
            - text: Extracted text
            - metadata: File metadata
            - content_type: Detected content type
            - file_size: File size in bytes
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        # Get file info
        file_size = file_path.stat().st_size
        content_type = content_type or self._magic.from_file(str(file_path), mime=True)
        file_extension = self._get_file_extension(file_path, content_type)
        
        # Initialize result
        result = {
            'text': '',
            'metadata': {
                'file_name': file_path.name,
                'file_size': file_size,
                'content_type': content_type,
                'file_extension': file_extension,
                'processing_time': datetime.utcnow().isoformat()
            },
            'content_type': content_type,
            'file_size': file_size
        }
        
        try:
            # Handle different file types
            if 'pdf' in content_type.lower():
                result['text'] = self._extract_text_from_pdf(file_path)
            elif 'wordprocessingml' in content_type.lower():
                result['text'] = self._extract_text_from_docx(file_path)
            elif 'presentationml' in content_type.lower():
                result['text'] = self._extract_text_from_pptx(file_path)
            elif 'spreadsheetml' in content_type.lower() or 'excel' in content_type.lower():
                result['text'] = self._extract_text_from_xlsx(file_path)
            elif 'text/' in content_type:
                result['text'] = file_path.read_text(encoding='utf-8', errors='ignore')
            elif any(img_type in content_type.lower() for img_type in ['image/png', 'image/jpeg', 'image/tiff', 'image/bmp']):
                result['text'] = self._extract_text_with_ocr(file_path)
            else:
                logger.warning(f"Unsupported content type for text extraction: {content_type}")
                
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            result['error'] = str(e)
            
            # Fallback to OCR if available and file is an image
            if self._is_ocr_supported(file_extension):
                try:
                    result['text'] = self._extract_text_with_ocr(file_path)
                    result['metadata']['ocr_used'] = True
                except Exception as ocr_error:
                    result['ocr_error'] = str(ocr_error)
        
        return result
    
    async def fetch_page(self, url: str, wait_for_selector: str = None, timeout: int = None) -> str:
        """
        Fetch a web page and return its content.
        
        Args:
            url: The URL to fetch
            wait_for_selector: Optional CSS selector to wait for before getting content
            timeout: Maximum time to wait in milliseconds
            
        Returns:
            str: The page content as HTML
        """
        if timeout is None:
            timeout = settings.BROWSER_TIMEOUT * 1000  # Convert to milliseconds
            
        try:
            logger.info(f"Fetching URL: {url}")
            await self.page.goto(url, timeout=timeout, wait_until="domcontentloaded")
            
            # Wait for the specified selector if provided
            if wait_for_selector:
                logger.debug(f"Waiting for selector: {wait_for_selector}")
                await self.page.wait_for_selector(wait_for_selector, timeout=timeout)
            
            # Get the page content
            content = await self.page.content()
            return content
            
        except Exception as e:
            logger.error(f"Error fetching page {url}: {str(e)}")
            raise
    
    async def download_file(self, url: str, save_path: Optional[Path] = None) -> Dict[str, Any]:
        """
        Download a file from a given URL.
        
        Args:
            url: URL of the file to download
            save_path: Optional path to save the file. If None, uses settings.DOWNLOADS_DIR
            
        Returns:
            Dict containing file metadata and content
        """
        try:
            if save_path is None:
                save_path = Path(settings.DOWNLOADS_DIR) / Path(url).name
            else:
                save_path = Path(save_path)
                
            save_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Start waiting for the download
            async with self.page.expect_download() as download_info:
                await self.page.goto(url)
            
            download = await download_info.value
            
            # Save the file
            await download.save_as(str(save_path))
            
            # Get file metadata
            file_type = magic.from_file(str(save_path), mime=True)
            file_size = save_path.stat().st_size
            
            # Read file content based on type
            with open(save_path, 'rb') as f:
                if file_type.startswith('text/'):
                    content = f.read().decode('utf-8')
                elif file_type.startswith('image/'):
                    content = f"data:{file_type};base64," + base64.b64encode(f.read()).decode('utf-8')
                else:
                    content = None
            
            return {
                'url': url,
                'path': str(save_path),
                'type': file_type,
                'size': file_size,
                'content': content
            }
            
        except Exception as e:
            logger.error(f"Error downloading file from {url}: {str(e)}")
            raise
    
    async def submit_quiz_answers(self, url: str, answers: List[Dict[str, Any]], max_retries: int = 3) -> Dict[str, Any]:
        """
        Submit quiz answers to the specified URL.
        
        Args:
            url: The quiz submission URL
            answers: List of answer dictionaries
            max_retries: Maximum number of retry attempts
            
        Returns:
            Dict containing submission result
        """
        retry_count = 0
        last_error = None
        
        while retry_count < max_retries:
            try:
                await self.page.goto(url, wait_until='networkidle')
                
                for answer in answers:
                    question_id = answer.get('question_id')
                    answer_value = answer.get('answer')
                    
                    if not question_id or answer_value is None:
                        continue
                        
                    # Handle different answer types
                    if isinstance(answer_value, (str, int, float, bool)):
                        # Text, number, or boolean answer
                        await self.page.fill(f'[name="{question_id}"]', str(answer_value))
                    elif isinstance(answer_value, dict) and 'file' in answer_value:
                        # File upload
                        file_path = answer_value.get('file')
                        if file_path and Path(file_path).exists():
                            file_input = await self.page.query_selector(f'[name="{question_id}"]')
                            if file_input:
                                await file_input.set_input_files(file_path)
                    
                # Submit the form
                submit_button = await self.page.query_selector('button[type="submit"], input[type="submit"]')
                if submit_button:
                    await submit_button.click()
                    await self.page.wait_for_load_state('networkidle')
                
                # Check for success/failure
                success = await self.page.evaluate('''() => {
                    return document.body.innerText.includes('success') || 
                           document.body.innerText.includes('thank you') ||
                           document.body.innerText.includes('submitted');
                }''')
                
                if success:
                    return {
                        'status': 'success',
                        'message': 'Quiz submitted successfully',
                        'url': self.page.url
                    }
                
                # If we got here, submission might have failed
                last_error = "Unknown submission error"
                
            except Exception as e:
                last_error = str(e)
                logger.warning(f"Attempt {retry_count + 1} failed: {last_error}")
                
            retry_count += 1
            if retry_count < max_retries:
                await asyncio.sleep(2 ** retry_count)  # Exponential backoff
        
        return {
            'status': 'error',
            'message': f'Failed to submit quiz after {max_retries} attempts',
            'error': last_error
        }
    
    async def extract_quiz_data(self, url: str) -> Dict[str, Any]:
        """
        Extract quiz data from a given URL, including questions, instructions,
        and submission information.
        
        Args:
            url: The quiz URL
            
        Returns:
            Dict containing quiz data with questions, instructions, and metadata
        """
        try:
            logger.info(f"Extracting quiz data from {url}")
            await self.page.goto(url, wait_until='domcontentloaded')
            
            # Wait for any dynamic content to load
            await asyncio.sleep(2)  # Give some time for JS to execute
            
            # Extract the entire page content first
            content = await self.page.content()
            
            # Extract metadata and questions using JavaScript evaluation
            quiz_data = await self.page.evaluate('''() => {
                // Try to find the main content container
                const mainContent = document.querySelector('main, .quiz-container, .content, body');
                
                // Extract the visible text content
                const getVisibleText = (element) => {
                    const text = [];
                    const walker = document.createTreeWalker(
                        element || document.body,
                        NodeFilter.SHOW_TEXT,
                        null,
                        false
                    );
                    
                    let node;
                    while (node = walker.nextNode()) {
                        if (node.parentElement.tagName === 'SCRIPT' || 
                            node.parentElement.tagName === 'STYLE' ||
                            node.parentElement.hidden ||
                            window.getComputedStyle(node.parentElement).display === 'none') {
                            continue;
                        }
                        text.push(node.textContent.trim());
                    }
                    
                    return text.filter(t => t.length > 0).join('\n');
                };
                
                // Extract all links that might be important
                const links = Array.from(document.querySelectorAll('a[href]')).map(a => ({
                    text: a.textContent.trim(),
                    href: a.href,
                    is_download: a.download !== '' || 
                               a.href.toLowerCase().endsWith('.pdf') ||
                               a.href.toLowerCase().match(/\.(pdf|docx?|xlsx?|csv|json|txt)$/i)
                }));
                
                // Find potential submission forms
                const forms = Array.from(document.forms).map(form => ({
                    action: form.action || window.location.href,
                    method: form.method || 'GET',
                    inputs: Array.from(form.elements).map(el => ({
                        name: el.name,
                        type: el.type,
                        value: el.value,
                        required: el.required,
                        tagName: el.tagName
                    }))
                }));
                
                // Extract any visible text that might contain instructions
                const visibleText = getVisibleText(mainContent);
                
                return {
                    title: document.title,
                    url: window.location.href,
                    content: visibleText,
                    links: links.filter(link => link.text.length > 0),
                    forms: forms,
                    metadata: {
                        extracted_at: new Date().toISOString(),
                        has_forms: forms.length > 0,
                        has_downloads: links.some(link => link.is_download)
                    },
                    raw_html: document.documentElement.outerHTML
                };
            }''')
            
            # Process the extracted data further in Python
            quiz_data['instructions'] = self._extract_instructions(quiz_data.get('content', ''))
            quiz_data['questions'] = self._identify_questions(quiz_data)
            
            # Add metadata about the extraction
            quiz_data['metadata']['extraction_method'] = 'automated'
            quiz_data['metadata']['processing_time'] = str(datetime.utcnow())
            
            logger.info(f"Successfully extracted quiz data from {url}")
            return quiz_data
            
        except Exception as e:
            logger.error(f"Error extracting quiz data from {url}: {str(e)}", exc_info=True)
            raise
            
    def _extract_instructions(self, content: str) -> List[Dict[str, Any]]:
        """Extract instructions from the quiz content."""
        instructions = []
        
        # Look for common instruction patterns
        instruction_patterns = [
            (r'(?i)instructions?:\s*(.*?)(?=\n\s*\d+\.|\n\s*[A-Z]|\.\s*\n|$)', 'general'),
            (r'(?i)note[s]?:\s*(.*?)(?=\n\s*\d+\.|\n\s*[A-Z]|\.\s*\n|$)', 'note'),
            (r'(?i)important:\s*(.*?)(?=\n\s*\d+\.|\n\s*[A-Z]|\.\s*\n|$)', 'important'),
            (r'(?i)hint:\s*(.*?)(?=\n\s*\d+\.|\n\s*[A-Z]|\.\s*\n|$)', 'hint')
        ]
        
        for pattern, instr_type in instruction_patterns:
            matches = re.finditer(pattern, content, re.DOTALL)
            for match in matches:
                instruction_text = match.group(1).strip()
                if instruction_text and len(instruction_text) > 10:  # Filter out very short matches
                    instructions.append({
                        'type': instr_type,
                        'text': instruction_text,
                        'context': content[max(0, match.start() - 50):min(len(content), match.end() + 50)]
                    })
        
        return instructions
    
    def _identify_questions(self, quiz_data: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Identify and extract questions from the quiz data."""
        questions = []
        content = quiz_data.get('content', '')
        
        # Simple regex pattern to identify questions (can be enhanced)
        question_patterns = [
            (r'(?:\n|^)\s*(\d+[\.\)])\s*(.*?)(?=\n\s*\d+[\.\)]|\n\s*[A-Z]|\.\s*\n|$)', 'numbering'),
            (r'(?:\n|^)\s*([A-Z][^\n?]+\?)(?=\s|$)', 'question_mark'),
            (r'(?:\n|^)\s*(Q\d+[:\.]?\s*)(.*?)(?=\n\s*Q\d|\n\s*[A-Z]|$)', 'q_prefix')
        ]
        
        for pattern, q_type in question_patterns:
            matches = list(re.finditer(pattern, content, re.DOTALL | re.IGNORECASE))
            for i, match in enumerate(matches):
                question_text = match.group(2 if len(match.groups()) > 1 else 1).strip()
                if len(question_text) < 5:  # Skip very short questions
                    continue
                    
                # Get the context around the question
                start = max(0, match.start() - 100)
                end = min(len(content), match.end() + 500)  # Get more content after the question
                context = content[start:end]
                
                questions.append({
                    'id': f"q{len(questions) + 1}",
                    'type': q_type,
                    'text': question_text,
                    'context': context,
                    'position': match.start(),
                    'metadata': {
                        'detected_by': pattern,
                        'context_length': len(context)
                    }
                })
        
        return questions
    
    async def submit_quiz_answer(self, url: str, answer_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Submit an answer to a quiz.
        
        Args:
            url: The quiz submission URL
            answer_data: Dictionary containing the answer and metadata
                {
                    'email': str,
                    'secret': str,
                    'answer': Any,
                    'metadata': Dict[str, Any]  # Optional additional data
                }
                
        Returns:
            Dict containing submission result
        """
        try:
            logger.info(f"Submitting answer to {url}")
            
            # Navigate to the submission URL
            await self.page.goto(url, wait_until='domcontentloaded')
            
            # Check if this is a form submission
            forms = await self.page.query_selector_all('form')
            
            if forms:
                # Handle form submission
                result = await self._submit_via_form(forms[0], answer_data)
            else:
                # Handle API submission
                result = await self._submit_via_api(url, answer_data)
                
            logger.info(f"Answer submitted successfully to {url}")
            return result
            
        except Exception as e:
            logger.error(f"Error submitting answer to {url}: {str(e)}", exc_info=True)
            raise
    
    async def _submit_via_form(self, form: ElementHandle, answer_data: Dict[str, Any]) -> Dict[str, Any]:
        """Submit answers via HTML form."""
        # Fill in the form fields
        for field, value in answer_data.items():
            if field in ['email', 'secret', 'answer']:
                # Handle different input types
                input_field = await form.query_selector(f'[name="{field}"], [id="{field}"]')
                if input_field:
                    input_type = await input_field.get_attribute('type')
                    
                    if input_type in ['text', 'email', 'hidden']:
                        await input_field.fill(str(value))
                    elif input_type in ['radio', 'checkbox']:
                        await input_field.check()
                    elif input_type == 'file':
                        # Handle file uploads if needed
                        if isinstance(value, str) and Path(value).exists():
                            await input_field.set_input_files(value)
        
        # Submit the form
        await form.evaluate('form => form.submit()')
        await self.page.wait_for_load_state('networkidle')
        
        # Get the response
        content = await self.page.content()
        return {
            'success': True,
            'submitted_at': datetime.utcnow().isoformat(),
            'response': content[:1000]  # Return first 1000 chars of response
        }
    
    async def _submit_via_api(self, url: str, answer_data: Dict[str, Any]) -> Dict[str, Any]:
        """Submit answers via API endpoint."""
        # Make a POST request to the API endpoint
        response = await self.page.request.post(
            url,
            headers={
                'Content-Type': 'application/json',
                'Accept': 'application/json'
            },
            data=json.dumps(answer_data)
        )
        
        # Parse the response
        try:
            result = await response.json()
        except:
            result = {'raw_response': await response.text()}
            
        return {
            'success': 200 <= response.status < 300,
            'status_code': response.status,
            'response': result,
            'submitted_at': datetime.utcnow().isoformat()
        }

# Example usage:
# async def example():
#     async with BrowserManager() as browser:
#         data = await browser.extract_quiz_data("https://example.com/quiz")
#         print(data)
# 
# if __name__ == "__main__":
#     import asyncio
#     asyncio.run(example())
